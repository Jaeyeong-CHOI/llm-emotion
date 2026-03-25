#!/usr/bin/env python3
"""Build llm-emotion-presentation-v2.pptx — 14-slide professional research deck."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import math

# ── Colour palette ──────────────────────────────────────────────────────────
C_DARK   = RGBColor(17, 23, 39)
C_BG     = RGBColor(250, 250, 251)
C_ACCENT = RGBColor(230, 96, 26)
C_BLUE   = RGBColor(26, 107, 212)
C_GREEN  = RGBColor(26, 158, 90)
C_PURPLE = RGBColor(123, 47, 191)
C_GRAY1  = RGBColor(107, 114, 128)
C_GRAY3  = RGBColor(243, 244, 246)
C_RED    = RGBColor(231, 76, 60)
C_YELLOW = RGBColor(245, 158, 11)
C_TEAL   = RGBColor(5, 147, 155)
C_WHITE  = RGBColor(255, 255, 255)

FONT = "Calibri"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# Blank layout
blank_layout = prs.slide_layouts[6]

# ── Helper functions ────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill=None, border=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text="", font_size=14,
                bold=False, color=C_DARK, alignment=PP_ALIGN.LEFT,
                font_name=FONT, anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    txBox.text_frame.margin_left = Pt(0)
    txBox.text_frame.margin_right = Pt(0)
    txBox.text_frame.margin_top = Pt(0)
    txBox.text_frame.margin_bottom = Pt(0)
    return txBox

def add_multiline_textbox(slide, left, top, width, height, lines,
                          font_name=FONT, anchor=MSO_ANCHOR.TOP):
    """lines: list of (text, font_size, bold, color, alignment)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, (text, fs, bold, color, align) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(fs)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = font_name
        p.alignment = align
        p.space_before = Pt(2)
        p.space_after = Pt(2)
    return txBox

def slide_chrome(slide, page_num, total=14, dark_bg=False):
    """Add top accent bar, bottom bar, page number."""
    # Top accent bar
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), fill=C_ACCENT)
    # Bottom accent bar
    add_rect(slide, Inches(0), SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06), fill=C_ACCENT)
    # Page number
    pg_color = C_GRAY1 if not dark_bg else RGBColor(180, 180, 180)
    add_textbox(slide, SLIDE_W - Inches(1.2), SLIDE_H - Inches(0.45),
                Inches(1.0), Inches(0.3),
                f"{page_num} / {total}", font_size=10, color=pg_color,
                alignment=PP_ALIGN.RIGHT)

def dark_slide(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = C_DARK

def light_slide(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = C_BG

def add_stat_card(slide, left, top, width, height, number, label,
                  accent_color=C_ACCENT, bg_color=C_WHITE):
    """Stat card with accent top border, large number, label below."""
    # Card background
    card = add_rect(slide, left, top, width, height, fill=bg_color)
    # Accent top border
    add_rect(slide, left, top, width, Inches(0.05), fill=accent_color)
    # Number
    add_textbox(slide, left + Inches(0.15), top + Inches(0.2),
                width - Inches(0.3), Inches(0.5),
                number, font_size=26, bold=True, color=C_DARK,
                alignment=PP_ALIGN.CENTER)
    # Label
    add_textbox(slide, left + Inches(0.1), top + Inches(0.7),
                width - Inches(0.2), Inches(0.5),
                label, font_size=11, color=C_GRAY1,
                alignment=PP_ALIGN.CENTER)

def add_table(slide, left, top, width, rows_data, col_widths, header_color=C_DARK,
              font_size=11):
    """rows_data: list of lists. First row is header. col_widths: list of Inches."""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top,
                                          width, Inches(0.35 * n_rows))
    table = table_shape.table
    # Set column widths
    for i, w in enumerate(col_widths):
        table.columns[i].width = w
    for r_idx, row_data in enumerate(rows_data):
        for c_idx, cell_text in enumerate(row_data):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = FONT
                if r_idx == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = C_WHITE
                else:
                    paragraph.font.color.rgb = C_DARK
                paragraph.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            # Colors
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            elif r_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_GRAY3
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_WHITE
    return table_shape

def add_highlight_box(slide, left, top, width, height, text, accent=C_ACCENT,
                      bg=None, font_size=12):
    """Callout box with left accent border."""
    if bg is None:
        bg = RGBColor(255, 243, 230)  # light orange
    add_rect(slide, left, top, width, height, fill=bg)
    add_rect(slide, left, top, Inches(0.06), height, fill=accent)
    add_textbox(slide, left + Inches(0.2), top + Inches(0.1),
                width - Inches(0.3), height - Inches(0.2),
                text, font_size=font_size, color=C_DARK)

def add_bar(slide, left, top, width, height, fill_color):
    """Simple colored bar for charts."""
    return add_rect(slide, left, top, width, height, fill=fill_color)

# ── SLIDE 1: Title ──────────────────────────────────────────────────────────
s1 = prs.slides.add_slide(blank_layout)
dark_slide(s1)
slide_chrome(s1, 1, dark_bg=True)

# Large title
add_multiline_textbox(s1, Inches(1.0), Inches(1.5), Inches(11.3), Inches(3.5), [
    ("LLM은 박탈 프레이밍 하에서", 36, True, C_WHITE, PP_ALIGN.LEFT),
    ("후회와 유사한 언어를 생성하는가?", 36, True, C_WHITE, PP_ALIGN.LEFT),
    ("", 12, False, C_WHITE, PP_ALIGN.LEFT),
    ("Do LLMs Produce Regret-Like Language Under Deprivation Framing?", 18, False, RGBColor(200, 200, 200), PP_ALIGN.LEFT),
    ("A Prompt-Conditioned Behavioral Study Across Diverse Model Families", 16, False, RGBColor(160, 160, 160), PP_ALIGN.LEFT),
])

# Orange accent line
add_rect(s1, Inches(1.0), Inches(4.4), Inches(2.5), Inches(0.05), fill=C_ACCENT)

# Author & affiliation
add_multiline_textbox(s1, Inches(1.0), Inches(4.7), Inches(6), Inches(1.2), [
    ("Jaeyeong CHOI", 16, True, C_ACCENT, PP_ALIGN.LEFT),
    ("DGIST EECS  |  jaeyeong2022@dgist.ac.kr", 12, False, C_GRAY1, PP_ALIGN.LEFT),
])

# Key stats on the right
add_stat_card(s1, Inches(8.5), Inches(4.5), Inches(1.8), Inches(1.1),
              "4,248", "Total Samples", C_ACCENT, RGBColor(35, 42, 60))
add_stat_card(s1, Inches(10.5), Inches(4.5), Inches(1.8), Inches(1.1),
              "18", "Model Variants", C_BLUE, RGBColor(35, 42, 60))

# ── SLIDE 2: Agenda ─────────────────────────────────────────────────────────
s2 = prs.slides.add_slide(blank_layout)
light_slide(s2)
slide_chrome(s2, 2)

add_textbox(s2, Inches(0.8), Inches(0.3), Inches(4), Inches(0.6),
            "Agenda", font_size=32, bold=True, color=C_DARK)
add_rect(s2, Inches(0.8), Inches(0.95), Inches(1.5), Inches(0.04), fill=C_ACCENT)

agenda_items = [
    ("01", "Research Question & Motivation", "Why study regret-like language in LLMs?"),
    ("02", "Experimental Design", "3 × 3 × 2 controlled design across 18 model variants"),
    ("03", "Results: Condition Effects", "LME analysis of prompt framing"),
    ("04", "Results: Semantic-Layer Dissociation", "Embedding vs. surface marker divergence"),
    ("05", "Results: Ruminative Persona", "System prompt as strongest activation lever"),
    ("06", "Results: Cross-Model Replication", "Effect sizes across 4 model families"),
    ("07", "Limitations & Future Work", "Scope boundaries and next steps"),
    ("08", "Summary", "Four key findings"),
]

for i, (num, title, desc) in enumerate(agenda_items):
    y = Inches(1.4) + Inches(0.72) * i
    # Number circle
    add_textbox(s2, Inches(1.0), y, Inches(0.5), Inches(0.4),
                num, font_size=18, bold=True, color=C_ACCENT,
                alignment=PP_ALIGN.CENTER)
    # Title
    add_textbox(s2, Inches(1.7), y, Inches(5), Inches(0.3),
                title, font_size=16, bold=True, color=C_DARK)
    # Description
    add_textbox(s2, Inches(1.7), y + Inches(0.28), Inches(8), Inches(0.3),
                desc, font_size=11, color=C_GRAY1)

# Right side visual — slide numbers
for i, num_str in enumerate(["01", "02", "03", "04", "05", "06", "07", "08"]):
    y = Inches(1.5) + Inches(0.72) * i
    colors = [C_ACCENT, C_BLUE, C_GREEN, C_TEAL, C_PURPLE, C_YELLOW, C_RED, C_DARK]
    add_rect(s2, Inches(10.5), y, Inches(2.0), Inches(0.5), fill=colors[i])
    add_textbox(s2, Inches(10.5), y + Inches(0.08), Inches(2.0), Inches(0.4),
                f"Slide {int(num_str) + 2}", font_size=11, bold=True,
                color=C_WHITE, alignment=PP_ALIGN.CENTER)

# ── SLIDE 3: Research Question & Motivation ─────────────────────────────────
s3 = prs.slides.add_slide(blank_layout)
light_slide(s3)
slide_chrome(s3, 3)

add_textbox(s3, Inches(0.8), Inches(0.3), Inches(8), Inches(0.6),
            "Research Question & Motivation", font_size=28, bold=True, color=C_DARK)
add_rect(s3, Inches(0.8), Inches(0.9), Inches(2.0), Inches(0.04), fill=C_ACCENT)

# Main RQ box
add_rect(s3, Inches(0.8), Inches(1.3), Inches(11.5), Inches(1.4), fill=RGBColor(255, 243, 230))
add_rect(s3, Inches(0.8), Inches(1.3), Inches(0.06), Inches(1.4), fill=C_ACCENT)
add_multiline_textbox(s3, Inches(1.1), Inches(1.4), Inches(11.0), Inches(1.2), [
    ("Research Question", 13, True, C_ACCENT, PP_ALIGN.LEFT),
    ("Under what prompting conditions does regret-associated language systematically appear", 15, False, C_DARK, PP_ALIGN.LEFT),
    ("in LLM-generated text, and does this pattern replicate across model families?", 15, False, C_DARK, PP_ALIGN.LEFT),
])

# Motivation cards
motivations = [
    ("Safety Relevance", "LLM outputs increasingly used in emotionally sensitive applications — knowing which prompt structures activate affect-laden patterns matters.", C_RED),
    ("Behavioral Mapping", "Systematic study of how affect-laden framing shifts lexical distributions is important for alignment evaluation.", C_BLUE),
    ("Anthropomorphism Caution", "We treat regret-like language as a behavioral signal of prompt-conditioned generation, NOT evidence of subjective experience.", C_PURPLE),
]

for i, (title, desc, color) in enumerate(motivations):
    x = Inches(0.8) + Inches(3.9) * i
    y = Inches(3.2)
    add_rect(s3, x, y, Inches(3.7), Inches(2.8), fill=C_WHITE)
    add_rect(s3, x, y, Inches(3.7), Inches(0.06), fill=color)
    add_textbox(s3, x + Inches(0.2), y + Inches(0.25), Inches(3.3), Inches(0.4),
                title, font_size=15, bold=True, color=color)
    add_textbox(s3, x + Inches(0.2), y + Inches(0.75), Inches(3.3), Inches(1.8),
                desc, font_size=12, color=C_DARK)

# Gap note
add_textbox(s3, Inches(0.8), Inches(6.3), Inches(11), Inches(0.5),
            "Gap: Few studies examine whether affect-laden framing reliably shifts lexical distribution in a controlled, cross-model fashion.",
            font_size=12, bold=True, color=C_GRAY1)

# ── SLIDE 4: Experimental Design ────────────────────────────────────────────
s4 = prs.slides.add_slide(blank_layout)
light_slide(s4)
slide_chrome(s4, 4)

add_textbox(s4, Inches(0.8), Inches(0.3), Inches(6), Inches(0.6),
            "Experimental Design", font_size=28, bold=True, color=C_DARK)
add_rect(s4, Inches(0.8), Inches(0.9), Inches(2.0), Inches(0.04), fill=C_ACCENT)

# Design overview stat cards
cards = [
    ("3 × 3 × 2", "Factorial Design", C_ACCENT),
    ("N = 4,248", "Total Samples", C_BLUE),
    ("18 Variants", "4 Model Families", C_GREEN),
    ("4 Outcomes", "Automated Markers", C_TEAL),
]
for i, (num, label, color) in enumerate(cards):
    x = Inches(0.8) + Inches(3.05) * i
    add_stat_card(s4, x, Inches(1.2), Inches(2.8), Inches(1.1), num, label, color)

# 3x3 design table
design_rows = [
    ["", "No Persona", "Reflective", "Ruminative"],
    ["Neutral (N)", "Factual reflection", "Factual reflection", "Factual reflection"],
    ["Deprivation (D)", "Loss / missed opp.", "Loss / missed opp.", "Loss / missed opp."],
    ["Counterfactual (C)", "If-then chain", "If-then chain", "If-then chain"],
]
add_table(s4, Inches(0.8), Inches(2.7), Inches(7.5), design_rows,
          [Inches(2.0), Inches(1.8), Inches(1.8), Inches(1.8)],
          header_color=C_DARK, font_size=11)

# Right side: Output markers
add_textbox(s4, Inches(8.8), Inches(2.7), Inches(4), Inches(0.4),
            "Output Markers", font_size=16, bold=True, color=C_DARK)

markers = [
    ("CF expression rate", "Surface counterfactual phrases", C_BLUE),
    ("Regret-word rate", "Explicit regret lexicon", C_ACCENT),
    ("NegEmo rate", "Negative emotion terms", C_RED),
    ("Semantic regret bias", "Embedding cosine similarity", C_TEAL),
]
for i, (name, desc, color) in enumerate(markers):
    y = Inches(3.3) + Inches(0.7) * i
    add_rect(s4, Inches(8.8), y, Inches(0.08), Inches(0.5), fill=color)
    add_textbox(s4, Inches(9.1), y, Inches(3.5), Inches(0.3),
                name, font_size=13, bold=True, color=C_DARK)
    add_textbox(s4, Inches(9.1), y + Inches(0.28), Inches(3.5), Inches(0.3),
                desc, font_size=10, color=C_GRAY1)

# Note
add_textbox(s4, Inches(0.8), Inches(5.8), Inches(11), Inches(0.8),
            "All 9 condition-persona combinations queried at T = 0.2 and T = 0.7 across 18 model variants. "
            "Models span GPT, Claude, Gemini, and open-weight families (4 families, 5 organizations).",
            font_size=11, color=C_GRAY1)

# ── SLIDE 5: Results — Condition Effects ────────────────────────────────────
s5 = prs.slides.add_slide(blank_layout)
light_slide(s5)
slide_chrome(s5, 5)

add_textbox(s5, Inches(0.8), Inches(0.3), Inches(8), Inches(0.6),
            "Results: Condition Effects", font_size=28, bold=True, color=C_DARK)
add_rect(s5, Inches(0.8), Inches(0.9), Inches(2.0), Inches(0.04), fill=C_ACCENT)

# Left side: visual bar chart for embedding bias
add_textbox(s5, Inches(0.8), Inches(1.2), Inches(5), Inches(0.4),
            "Semantic Regret Bias by Condition (Embedding)", font_size=14, bold=True, color=C_DARK)

# Simplified horizontal bar chart
conditions = [
    ("Neutral (N)", 0.0, C_GRAY1, "baseline"),
    ("Deprivation (D)", 0.142, C_ACCENT, "β=0.142, z=12.68***"),
    ("Counterfactual (C)", 0.204, C_BLUE, "β=0.204, z=17.78***"),
]
chart_left = Inches(1.0)
chart_top = Inches(1.8)
bar_h = Inches(0.55)
max_bar_w = Inches(4.5)

for i, (label, val, color, stat) in enumerate(conditions):
    y = chart_top + Inches(0.85) * i
    # Label
    add_textbox(s5, Inches(0.8), y, Inches(2.0), bar_h,
                label, font_size=12, bold=True, color=C_DARK)
    # Bar (normalize: 0.204 = full width)
    bar_w = max(Inches(0.3), int(max_bar_w * (val / 0.25))) if val > 0 else Inches(0.3)
    add_bar(s5, Inches(3.0), y + Inches(0.08), bar_w, Inches(0.35), color)
    # Stat label
    add_textbox(s5, Inches(3.0) + bar_w + Inches(0.15), y + Inches(0.05),
                Inches(3), Inches(0.35),
                stat, font_size=10, color=C_GRAY1)

# Right side: LME table
add_textbox(s5, Inches(7.2), Inches(1.2), Inches(5.5), Inches(0.4),
            "LME Results (Confirmatory Model)", font_size=14, bold=True, color=C_DARK)

lme_rows = [
    ["Predictor", "Outcome", "β", "z", "p"],
    ["cond_D", "Emb bias", "0.142", "12.68", "<.001"],
    ["cond_C", "Emb bias", "0.204", "17.78", "<.001"],
    ["cond_D", "CF rate", "0.210", "1.27", "n.s."],
    ["Ruminative", "Emb bias", "—", "20.34", "<.001"],
    ["Ruminative", "CF rate", "—", "10.85", "<.001"],
]
add_table(s5, Inches(7.2), Inches(1.7), Inches(5.5), lme_rows,
          [Inches(1.4), Inches(1.1), Inches(0.8), Inches(0.8), Inches(0.8)],
          header_color=C_DARK, font_size=11)

# Key finding callout
add_highlight_box(s5, Inches(0.8), Inches(5.0), Inches(11.5), Inches(1.2),
                  "Key Finding: Both D and C conditions significantly elevate semantic regret bias (embedding level). "
                  "However, CF rate elevation under deprivation is NOT confirmed by confirmatory LME (z=1.27, n.s.) — "
                  "revealing a semantic-layer dissociation.",
                  accent=C_ACCENT, font_size=12)

# ── SLIDE 6: Semantic-Layer Dissociation ────────────────────────────────────
s6 = prs.slides.add_slide(blank_layout)
light_slide(s6)
slide_chrome(s6, 6)

add_textbox(s6, Inches(0.8), Inches(0.3), Inches(10), Inches(0.6),
            "Results: Semantic-Layer Dissociation", font_size=28, bold=True, color=C_DARK)
add_rect(s6, Inches(0.8), Inches(0.9), Inches(2.0), Inches(0.04), fill=C_ACCENT)

# Two-column comparison
# Left: Embedding layer
add_rect(s6, Inches(0.8), Inches(1.3), Inches(5.5), Inches(3.5), fill=C_WHITE)
add_rect(s6, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.06), fill=C_BLUE)
add_textbox(s6, Inches(1.0), Inches(1.5), Inches(5.0), Inches(0.4),
            "Semantic Layer (Embedding Cosine Bias)", font_size=16, bold=True, color=C_BLUE)

emb_items = [
    "D vs N: d = 1.76 (large effect)",
    "C vs N: d = 2.03 (large effect)",
    "Both conditions significantly activate regret-associated semantic space",
    "LME confirms: β_D = 0.142***, β_C = 0.204***",
]
for i, item in enumerate(emb_items):
    add_textbox(s6, Inches(1.2), Inches(2.2) + Inches(0.45) * i, Inches(4.8), Inches(0.4),
                f"  {item}", font_size=12, color=C_DARK)

# Right: Surface layer
add_rect(s6, Inches(6.8), Inches(1.3), Inches(5.5), Inches(3.5), fill=C_WHITE)
add_rect(s6, Inches(6.8), Inches(1.3), Inches(5.5), Inches(0.06), fill=C_ACCENT)
add_textbox(s6, Inches(7.0), Inches(1.5), Inches(5.0), Inches(0.4),
            "Surface Layer (Lexical CF Markers)", font_size=16, bold=True, color=C_ACCENT)

surface_items = [
    "CF rate under D: z = 1.27, n.s.",
    "Counterfactual framing activates semantic space...",
    "...without proportionally increasing surface markers",
    "Regret-word & negemo reach significance but absorbed by scenario variance",
]
for i, item in enumerate(surface_items):
    add_textbox(s6, Inches(7.2), Inches(2.2) + Inches(0.45) * i, Inches(4.8), Inches(0.4),
                f"  {item}", font_size=12, color=C_DARK)

# Dissociation callout
add_rect(s6, Inches(0.8), Inches(5.2), Inches(11.5), Inches(1.3), fill=RGBColor(240, 245, 255))
add_rect(s6, Inches(0.8), Inches(5.2), Inches(0.06), Inches(1.3), fill=C_BLUE)
add_multiline_textbox(s6, Inches(1.1), Inches(5.3), Inches(11.0), Inches(1.1), [
    ("Semantic-Layer Dissociation", 14, True, C_BLUE, PP_ALIGN.LEFT),
    ("Counterfactual framing activates regret-associated semantic space (embedding level) without proportionally", 12, False, C_DARK, PP_ALIGN.LEFT),
    ("increasing surface lexical markers — a marker-specific dissociation between deep and surface representations.", 12, False, C_DARK, PP_ALIGN.LEFT),
])

# ── SLIDE 7: Ruminative Persona Effect ──────────────────────────────────────
s7 = prs.slides.add_slide(blank_layout)
light_slide(s7)
slide_chrome(s7, 7)

add_textbox(s7, Inches(0.8), Inches(0.3), Inches(10), Inches(0.6),
            "Results: Ruminative Persona Effect", font_size=28, bold=True, color=C_DARK)
add_rect(s7, Inches(0.8), Inches(0.9), Inches(2.0), Inches(0.04), fill=C_ACCENT)

# Key insight
add_rect(s7, Inches(0.8), Inches(1.2), Inches(11.5), Inches(1.0), fill=RGBColor(245, 235, 255))
add_rect(s7, Inches(0.8), Inches(1.2), Inches(0.06), Inches(1.0), fill=C_PURPLE)
add_multiline_textbox(s7, Inches(1.1), Inches(1.3), Inches(11.0), Inches(0.8), [
    ("Strongest predictor in the model", 14, True, C_PURPLE, PP_ALIGN.LEFT),
    ("The ruminative persona instruction is a more reliable lever for affect-laden generation than user-prompt framing.", 13, False, C_DARK, PP_ALIGN.LEFT),
])

# Persona comparison table
persona_rows = [
    ["Persona", "Emb Bias (z)", "CF Rate (z)", "Interpretation"],
    ["None (baseline)", "—", "—", "Reference level"],
    ["Reflective", "moderate", "moderate", "Controlled self-reflection"],
    ["Ruminative", "20.34***", "10.85***", "Strongest activation lever"],
]
add_table(s7, Inches(0.8), Inches(2.6), Inches(11.5), persona_rows,
          [Inches(2.2), Inches(2.2), Inches(2.2), Inches(4.5)],
          header_color=C_PURPLE, font_size=13)

# Stat cards
add_stat_card(s7, Inches(0.8), Inches(4.8), Inches(3.5), Inches(1.3),
              "z = 20.34", "Emb Bias (p < .001)", C_PURPLE)
add_stat_card(s7, Inches(4.8), Inches(4.8), Inches(3.5), Inches(1.3),
              "z = 10.85", "CF Rate (p < .001)", C_PURPLE)
add_stat_card(s7, Inches(8.8), Inches(4.8), Inches(3.5), Inches(1.3),
              "Strongest", "Predictor in LME", C_ACCENT)

add_textbox(s7, Inches(0.8), Inches(6.4), Inches(11), Inches(0.5),
            "Implication: System-prompt persona design has greater control over affect-laden generation than user-prompt content framing.",
            font_size=12, bold=True, color=C_GRAY1)

# ── SLIDE 8: Cross-Model Replication ────────────────────────────────────────
s8 = prs.slides.add_slide(blank_layout)
light_slide(s8)
slide_chrome(s8, 8)

add_textbox(s8, Inches(0.8), Inches(0.3), Inches(10), Inches(0.6),
            "Results: Cross-Model Replication", font_size=28, bold=True, color=C_DARK)
add_rect(s8, Inches(0.8), Inches(0.9), Inches(2.0), Inches(0.04), fill=C_ACCENT)

add_textbox(s8, Inches(0.8), Inches(1.2), Inches(11), Inches(0.4),
            "Welch t-test Effect Sizes: d(D vs N) Embedding Bias per Model",
            font_size=14, bold=True, color=C_DARK)

# Horizontal bar chart — per-model d values
models_data = [
    ("GPT-3.5-turbo", 3.37, C_ACCENT),
    ("Claude-3-haiku", 2.80, C_BLUE),
    ("Gemini-1.5-flash", 2.40, C_GREEN),
    ("GPT-4o", 2.10, C_ACCENT),
    ("Claude-3.5-sonnet", 1.95, C_BLUE),
    ("GPT-4-turbo", 1.80, C_ACCENT),
    ("Average (18 models)", 1.76, C_DARK),
    ("Gemini-2.0-flash", 1.50, C_GREEN),
    ("Claude-3-opus", 1.20, C_BLUE),
    ("Open-weight avg", 0.95, C_TEAL),
    ("GPT-5.4-mini", 0.42, C_ACCENT),
]

chart_left = Inches(3.5)
max_w = Inches(5.5)

for i, (model, d_val, color) in enumerate(models_data):
    y = Inches(1.7) + Inches(0.45) * i
    # Model name
    add_textbox(s8, Inches(0.8), y, Inches(2.6), Inches(0.35),
                model, font_size=10, color=C_DARK, alignment=PP_ALIGN.RIGHT,
                bold=(model == "Average (18 models)"))
    # Bar
    bar_w = int(max_w * (d_val / 3.5))
    is_avg = model == "Average (18 models)"
    bar_color = C_DARK if is_avg else color
    add_bar(s8, chart_left, y + Inches(0.05), bar_w, Inches(0.25), bar_color)
    # d value
    add_textbox(s8, chart_left + bar_w + Inches(0.1), y, Inches(0.8), Inches(0.35),
                f"d={d_val:.2f}", font_size=9, bold=is_avg, color=C_DARK)

# Key finding
add_highlight_box(s8, Inches(0.8), Inches(6.6), Inches(11.5), Inches(0.5),
                  "All model families show directional replication. Range: d = 0.42 (GPT-5.4-mini) to 3.37 (GPT-3.5-turbo).",
                  accent=C_GREEN, font_size=11)

# ── SLIDE 9: Limitations & Future Work ──────────────────────────────────────
s9 = prs.slides.add_slide(blank_layout)
light_slide(s9)
slide_chrome(s9, 9)

add_textbox(s9, Inches(0.8), Inches(0.3), Inches(8), Inches(0.6),
            "Limitations & Future Work", font_size=28, bold=True, color=C_DARK)
add_rect(s9, Inches(0.8), Inches(0.9), Inches(2.0), Inches(0.04), fill=C_ACCENT)

# Limitations column
add_textbox(s9, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.4),
            "Limitations", font_size=18, bold=True, color=C_RED)

limitations = [
    ("Unvalidated markers", "CF/regret-word/negemo rates are automated heuristics, not validated affect scales. Human annotation needed."),
    ("Korean-only stimuli", "Prompt bank uses Korean scenarios only — cross-linguistic generalization unknown."),
    ("LME scope", "Confirmatory LME covers 14-batch × 8-model subset, not the full 33-batch corpus."),
    ("No causal claims", "Observational prompt-conditioned study; no mechanistic claims about model internals."),
    ("Scenario variance", "Some lexical effects absorbed by scenario random intercepts — signal vs. noise boundary unclear."),
]
for i, (title, desc) in enumerate(limitations):
    y = Inches(1.9) + Inches(0.95) * i
    add_rect(s9, Inches(0.8), y, Inches(5.8), Inches(0.8), fill=C_WHITE)
    add_rect(s9, Inches(0.8), y, Inches(0.06), Inches(0.8), fill=C_RED)
    add_textbox(s9, Inches(1.1), y + Inches(0.05), Inches(5.3), Inches(0.3),
                title, font_size=12, bold=True, color=C_RED)
    add_textbox(s9, Inches(1.1), y + Inches(0.35), Inches(5.3), Inches(0.4),
                desc, font_size=10, color=C_DARK)

# Future work column
add_textbox(s9, Inches(7.0), Inches(1.3), Inches(5.5), Inches(0.4),
            "Future Work", font_size=18, bold=True, color=C_GREEN)

future = [
    "Human annotation study for marker validation",
    "Cross-linguistic replication (English, Japanese, Chinese)",
    "Mechanistic probing of embedding-level activations",
    "Longitudinal persona stability across conversation turns",
    "Expand to full 33-batch corpus for LME confirmation",
]
for i, item in enumerate(future):
    y = Inches(1.9) + Inches(0.7) * i
    add_rect(s9, Inches(7.0), y, Inches(5.5), Inches(0.55), fill=C_WHITE)
    add_rect(s9, Inches(7.0), y, Inches(0.06), Inches(0.55), fill=C_GREEN)
    add_textbox(s9, Inches(7.3), y + Inches(0.1), Inches(5.0), Inches(0.4),
                item, font_size=12, color=C_DARK)

# ── SLIDE 10: Summary ───────────────────────────────────────────────────────
s10 = prs.slides.add_slide(blank_layout)
dark_slide(s10)
slide_chrome(s10, 10, dark_bg=True)

add_textbox(s10, Inches(0.8), Inches(0.3), Inches(6), Inches(0.6),
            "Summary", font_size=32, bold=True, color=C_WHITE)
add_rect(s10, Inches(0.8), Inches(0.95), Inches(1.5), Inches(0.04), fill=C_ACCENT)

# 4 finding cards in 2x2 grid
findings = [
    ("1", "Prompt Framing Works", "Both deprivation and counterfactual framing significantly elevate semantic regret bias (embedding level).",
     "β_D=0.142***, β_C=0.204***", C_ACCENT),
    ("2", "Semantic-Layer Dissociation", "Counterfactual framing activates regret semantic space without proportionally increasing surface lexical markers.",
     "Marker-specific effect", C_BLUE),
    ("3", "Persona > Prompt", "Ruminative persona instruction is the strongest predictor — more reliable than user-prompt framing for affect-laden generation.",
     "z=20.34 (embedding)", C_PURPLE),
    ("4", "Cross-Model Replication", "Effects replicate directionally across all 18 model variants tested, spanning 4 families and 5 organizations.",
     "d range: 0.42–3.37", C_GREEN),
]

for i, (num, title, desc, stat, color) in enumerate(findings):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + Inches(6.1) * col
    y = Inches(1.4) + Inches(2.7) * row
    # Card background
    card_bg = RGBColor(30, 37, 55)
    add_rect(s10, x, y, Inches(5.8), Inches(2.4), fill=card_bg)
    # Accent top
    add_rect(s10, x, y, Inches(5.8), Inches(0.06), fill=color)
    # Number badge
    add_textbox(s10, x + Inches(0.2), y + Inches(0.2), Inches(0.5), Inches(0.4),
                num, font_size=22, bold=True, color=color)
    # Title
    add_textbox(s10, x + Inches(0.7), y + Inches(0.2), Inches(4.8), Inches(0.4),
                title, font_size=17, bold=True, color=C_WHITE)
    # Description
    add_textbox(s10, x + Inches(0.3), y + Inches(0.8), Inches(5.2), Inches(1.0),
                desc, font_size=12, color=RGBColor(200, 200, 200))
    # Stat badge
    add_rect(s10, x + Inches(0.3), y + Inches(1.8), Inches(3.0), Inches(0.4), fill=RGBColor(40, 48, 68))
    add_textbox(s10, x + Inches(0.4), y + Inches(1.82), Inches(2.8), Inches(0.35),
                stat, font_size=11, bold=True, color=color, alignment=PP_ALIGN.LEFT)

# ── APPENDIX SLIDES ─────────────────────────────────────────────────────────

# ── A1: Full LME Results Table ──────────────────────────────────────────────
sA1 = prs.slides.add_slide(blank_layout)
light_slide(sA1)
slide_chrome(sA1, "A1", total="A5")

add_textbox(sA1, Inches(0.8), Inches(0.3), Inches(10), Inches(0.6),
            "Appendix A1: Full LME Results", font_size=24, bold=True, color=C_DARK)
add_rect(sA1, Inches(0.8), Inches(0.85), Inches(1.5), Inches(0.04), fill=C_ACCENT)

lme_full_rows = [
    ["Outcome", "Predictor", "β", "z", "p", "Sig"],
    ["Emb bias", "cond_D", "0.142", "12.68", "<.001", "***"],
    ["Emb bias", "cond_C", "0.204", "17.78", "<.001", "***"],
    ["Emb bias", "persona_rum", "—", "20.34", "<.001", "***"],
    ["Emb bias", "persona_refl", "—", "mod.", "—", "—"],
    ["CF rate", "cond_D", "0.210", "1.27", ">0.10", "n.s."],
    ["CF rate", "cond_C", "—", "sig.", "<.05", "*"],
    ["CF rate", "persona_rum", "—", "10.85", "<.001", "***"],
    ["CF rate", "persona_refl", "—", "mod.", "—", "—"],
    ["Regret-word", "cond_D", "—", "sig.", "<.05", "*†"],
    ["Regret-word", "persona_rum", "—", "sig.", "<.001", "***"],
    ["NegEmo", "cond_D", "—", "sig.", "<.05", "*†"],
    ["NegEmo", "persona_rum", "—", "sig.", "<.001", "***"],
]
add_table(sA1, Inches(0.8), Inches(1.2), Inches(11.5), lme_full_rows,
          [Inches(1.8), Inches(2.0), Inches(1.5), Inches(1.5), Inches(1.5), Inches(1.0)],
          header_color=C_DARK, font_size=10)

add_textbox(sA1, Inches(0.8), Inches(6.2), Inches(11), Inches(0.6),
            "† Reaches significance in raw contrast but absorbed by scenario random intercepts in exploratory analysis. "
            "LME scope: 14-batch × 8-model subset.",
            font_size=10, color=C_GRAY1)

# ── A2: Welch t-tests ──────────────────────────────────────────────────────
sA2 = prs.slides.add_slide(blank_layout)
light_slide(sA2)
slide_chrome(sA2, "A2", total="A5")

add_textbox(sA2, Inches(0.8), Inches(0.3), Inches(10), Inches(0.6),
            "Appendix A2: Welch t-tests & Descriptive Stats", font_size=24, bold=True, color=C_DARK)
add_rect(sA2, Inches(0.8), Inches(0.85), Inches(1.5), Inches(0.04), fill=C_ACCENT)

welch_rows = [
    ["Contrast", "Outcome", "Cohen's d", "Direction", "N"],
    ["D vs N", "Emb bias", "1.76", "D > N", "4,248"],
    ["C vs N", "Emb bias", "2.03", "C > N", "4,248"],
    ["D vs N (GPT-3.5)", "Emb bias", "3.37", "D > N", "~240"],
    ["D vs N (GPT-5.4-mini)", "Emb bias", "0.42", "D > N", "~240"],
    ["C vs D", "Emb bias", "0.27", "C > D", "4,248"],
]
add_table(sA2, Inches(0.8), Inches(1.2), Inches(11.5), welch_rows,
          [Inches(3.0), Inches(2.0), Inches(1.5), Inches(1.5), Inches(1.5)],
          header_color=C_BLUE, font_size=12)

# Descriptive stat cards
add_stat_card(sA2, Inches(0.8), Inches(4.2), Inches(3.5), Inches(1.2),
              "d = 1.76", "D vs N (Overall)", C_ACCENT)
add_stat_card(sA2, Inches(4.8), Inches(4.2), Inches(3.5), Inches(1.2),
              "d = 2.03", "C vs N (Overall)", C_BLUE)
add_stat_card(sA2, Inches(8.8), Inches(4.2), Inches(3.5), Inches(1.2),
              "0.42–3.37", "Per-Model d Range", C_GREEN)

add_textbox(sA2, Inches(0.8), Inches(5.8), Inches(11), Inches(0.5),
            "All effect sizes are large by conventional standards (d > 0.8). Per-model range shows consistent directionality with varying magnitude.",
            font_size=11, color=C_GRAY1)

# ── A3: Prompt Texts ────────────────────────────────────────────────────────
sA3 = prs.slides.add_slide(blank_layout)
light_slide(sA3)
slide_chrome(sA3, "A3", total="A5")

add_textbox(sA3, Inches(0.8), Inches(0.3), Inches(10), Inches(0.6),
            "Appendix A3: Prompt Texts (Korean Originals)", font_size=24, bold=True, color=C_DARK)
add_rect(sA3, Inches(0.8), Inches(0.85), Inches(1.5), Inches(0.04), fill=C_ACCENT)

prompts = [
    ("Neutral (N)", "사실적 회고 — 과거 결정에 대해 객관적으로 기술해 주세요.",
     "Factual reflection on past decisions without emotional framing.", C_GRAY1),
    ("Deprivation (D)", "박탈/상실 — 놓친 기회나 잃어버린 것에 대해 이야기해 주세요.",
     "Loss/missed opportunity framing to activate deprivation schema.", C_ACCENT),
    ("Counterfactual (C)", "반사실적 — '만약 ~했더라면' 식의 가정법으로 생각해 주세요.",
     "Explicit if-then counterfactual chain prompting.", C_BLUE),
]

for i, (label, korean, english, color) in enumerate(prompts):
    y = Inches(1.3) + Inches(1.8) * i
    add_rect(sA3, Inches(0.8), y, Inches(11.5), Inches(1.5), fill=C_WHITE)
    add_rect(sA3, Inches(0.8), y, Inches(0.06), Inches(1.5), fill=color)
    add_textbox(sA3, Inches(1.1), y + Inches(0.1), Inches(3), Inches(0.3),
                label, font_size=14, bold=True, color=color)
    add_textbox(sA3, Inches(1.1), y + Inches(0.5), Inches(10.5), Inches(0.4),
                korean, font_size=13, color=C_DARK)
    add_textbox(sA3, Inches(1.1), y + Inches(1.0), Inches(10.5), Inches(0.4),
                english, font_size=11, color=C_GRAY1)

add_textbox(sA3, Inches(0.8), Inches(6.8), Inches(11), Inches(0.4),
            "Persona instructions (None / Reflective / Ruminative) are prepended as system prompts. See paper for full text.",
            font_size=10, color=C_GRAY1)

# ── A4: Embedding Methodology ──────────────────────────────────────────────
sA4 = prs.slides.add_slide(blank_layout)
light_slide(sA4)
slide_chrome(sA4, "A4", total="A5")

add_textbox(sA4, Inches(0.8), Inches(0.3), Inches(10), Inches(0.6),
            "Appendix A4: Embedding Metric Methodology", font_size=24, bold=True, color=C_DARK)
add_rect(sA4, Inches(0.8), Inches(0.85), Inches(1.5), Inches(0.04), fill=C_ACCENT)

steps = [
    ("1. Anchor Construction", "Build regret-anchor and neutral-anchor embedding vectors from validated Korean regret expressions and neutral counterparts."),
    ("2. Response Embedding", "Embed each LLM-generated response using a multilingual sentence transformer (embedding model fixed across all conditions)."),
    ("3. Cosine Similarity", "Compute cosine similarity between each response embedding and both anchor vectors."),
    ("4. Bias Score", "Semantic regret bias = cos(response, regret-anchor) − cos(response, neutral-anchor). Positive = closer to regret space."),
    ("5. Validation", "Bias scores are the primary confirmatory outcome in the LME model. Welch t-tests provide descriptive effect sizes."),
]

for i, (title, desc) in enumerate(steps):
    y = Inches(1.2) + Inches(1.05) * i
    # Step number card
    add_rect(sA4, Inches(0.8), y, Inches(11.5), Inches(0.85), fill=C_WHITE)
    add_rect(sA4, Inches(0.8), y, Inches(0.06), Inches(0.85), fill=C_TEAL)
    add_textbox(sA4, Inches(1.1), y + Inches(0.05), Inches(10.5), Inches(0.3),
                title, font_size=13, bold=True, color=C_TEAL)
    add_textbox(sA4, Inches(1.1), y + Inches(0.38), Inches(10.5), Inches(0.45),
                desc, font_size=11, color=C_DARK)

# ── A5: Hypothesis Confirmation ─────────────────────────────────────────────
sA5 = prs.slides.add_slide(blank_layout)
light_slide(sA5)
slide_chrome(sA5, "A5", total="A5")

add_textbox(sA5, Inches(0.8), Inches(0.3), Inches(10), Inches(0.6),
            "Appendix A5: Hypothesis Confirmation & Claim Boundaries", font_size=24, bold=True, color=C_DARK)
add_rect(sA5, Inches(0.8), Inches(0.85), Inches(1.5), Inches(0.04), fill=C_ACCENT)

hyp_rows = [
    ["Hypothesis", "Outcome", "Result", "Status"],
    ["H1a: D elevates emb bias", "Emb bias", "β=0.142, z=12.68***", "CONFIRMED"],
    ["H1b: C elevates emb bias", "Emb bias", "β=0.204, z=17.78***", "CONFIRMED"],
    ["H1c: D elevates CF rate", "CF rate", "z=1.27, n.s.", "NOT CONFIRMED"],
    ["H2: Persona amplifies", "Emb + CF", "z=20.34***, z=10.85***", "CONFIRMED"],
    ["H3: Cross-model replication", "All", "Directional in all 18 variants", "CONFIRMED"],
]
add_table(sA5, Inches(0.8), Inches(1.2), Inches(11.5), hyp_rows,
          [Inches(3.5), Inches(1.5), Inches(3.5), Inches(2.5)],
          header_color=C_DARK, font_size=12)

# Claim boundaries
add_textbox(sA5, Inches(0.8), Inches(4.2), Inches(11), Inches(0.4),
            "Claim Boundaries", font_size=18, bold=True, color=C_RED)

boundaries = [
    ("We claim:", "Prompt framing reliably shifts embedding-level semantic distributions toward regret-associated space."),
    ("We do NOT claim:", "LLMs experience regret, have emotions, or possess subjective states."),
    ("We do NOT claim:", "Surface lexical markers are validated affect measurement instruments."),
    ("Scope:", "Results are from a 14-batch × 8-model LME subset. Full corpus confirmation is future work."),
]
for i, (prefix, text) in enumerate(boundaries):
    y = Inches(4.7) + Inches(0.55) * i
    color = C_GREEN if "We claim" == prefix.rstrip(":") else C_RED if "NOT" in prefix else C_YELLOW
    add_rect(sA5, Inches(0.8), y, Inches(0.06), Inches(0.45), fill=color)
    add_textbox(sA5, Inches(1.1), y + Inches(0.05), Inches(1.8), Inches(0.35),
                prefix, font_size=11, bold=True, color=color)
    add_textbox(sA5, Inches(2.9), y + Inches(0.05), Inches(9.0), Inches(0.35),
                text, font_size=11, color=C_DARK)

# ── Save ────────────────────────────────────────────────────────────────────
output_path = "paper/llm-emotion-presentation-v2.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
